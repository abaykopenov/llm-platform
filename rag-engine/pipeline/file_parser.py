"""
File Parser — Extract text from various file formats.
Supports: PDF, DOCX, PPTX, XLSX, TXT, MD, CSV, HTML
"""

import csv
import io
import os
import re
from typing import List


def parse_txt(file_path: str) -> str:
    """Parse plain text or markdown file."""
    encodings = ["utf-8", "cp1251", "latin-1"]
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    with open(file_path, "rb") as f:
        return f.read().decode("utf-8", errors="replace")


def parse_pdf(file_path: str) -> str:
    """Parse PDF file using PyMuPDF (fitz)."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        pages = []
        for page in doc:
            text = page.get_text().strip()
            if text:
                pages.append(text)
        doc.close()
        return "\n\n".join(pages)
    except ImportError:
        # Fallback to PyPDF2
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text.strip())
        return "\n\n".join(pages)


def parse_docx(file_path: str) -> str:
    """Parse DOCX file using python-docx."""
    from docx import Document
    doc = Document(file_path)
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)
    return "\n\n".join(paragraphs)


def parse_pptx(file_path: str) -> str:
    """Parse PPTX file using python-pptx."""
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"--- Слайд {slide_num} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if len(slide_texts) > 1:
            texts.append("\n".join(slide_texts))
    return "\n\n".join(texts)


def parse_xlsx(file_path: str) -> str:
    """Parse XLSX file using openpyxl."""
    from openpyxl import load_workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    texts = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        texts.append(f"--- Лист: {sheet} ---")
        rows = []
        for row in ws.iter_rows(values_only=True):
            vals = [str(v) for v in row if v is not None]
            if vals:
                rows.append(" | ".join(vals))
        texts.append("\n".join(rows))
    wb.close()
    return "\n\n".join(texts)


def parse_csv(file_path: str) -> str:
    """Parse CSV file."""
    text_parts = []
    raw = parse_txt(file_path)
    reader = csv.reader(io.StringIO(raw))
    headers = None
    for i, row in enumerate(reader):
        if i == 0:
            headers = row
            continue
        if headers:
            line = "; ".join(f"{h}: {v}" for h, v in zip(headers, row) if v.strip())
        else:
            line = "; ".join(row)
        if line.strip():
            text_parts.append(line)
    return "\n".join(text_parts)


def parse_html(file_path: str) -> str:
    """Parse HTML file using BeautifulSoup."""
    from bs4 import BeautifulSoup
    raw = parse_txt(file_path)
    soup = BeautifulSoup(raw, "html.parser")
    # Remove script and style tags
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


# ─── Registry ────────────────────────────────────────────────

PARSERS = {
    ".txt": parse_txt,
    ".md": parse_txt,
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".xlsx": parse_xlsx,
    ".csv": parse_csv,
    ".html": parse_html,
    ".htm": parse_html,
}


def parse_file(file_path: str) -> str:
    """Parse a file based on its extension. Returns extracted text."""
    ext = os.path.splitext(file_path)[1].lower()
    parser = PARSERS.get(ext)
    if parser is None:
        supported = ", ".join(PARSERS.keys())
        raise ValueError(f"Unsupported file type: {ext}. Supported: {supported}")
    return parser(file_path)


def get_supported_extensions() -> List[str]:
    """Return list of supported file extensions."""
    return list(PARSERS.keys())
